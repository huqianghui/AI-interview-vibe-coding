"""Binary-format text extractors (coverage-omitted — need the ``azure`` extra's parsers).

pdfplumber / python-docx / python-pptx are optional deps not installed in CI, so these are
imported lazily and this module is coverage-omitted (same precedent as the azure adapters).
The dispatcher in ``extraction.py`` (CI-covered) wraps every call in try/except, so a missing
parser degrades to ``""`` rather than crashing ingestion.

Each extractor follows the reference ``skill_text_extractor`` pattern verbatim:
lazy import, try/except, return ``""`` on any failure (never raises).

Two shapes per format:
- ``extract_*`` returns one flattened string (kept for callers that only need text).
- ``segments_*`` returns ``[(label, text), ...]`` preserving the page/slide boundary so
  ingestion can persist a human-facing location label per chunk (SPEC F1 traceability, AC #1).
  A failure returns ``[]`` (the dispatcher degrades it the same way it degrades ``""``).
"""

import io
import logging

logger = logging.getLogger(__name__)


def extract_pdf(content: bytes) -> str:
    try:
        import pdfplumber

        pages: list[str] = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
        return "\n\n".join(pages)
    except Exception as exc:  # noqa: BLE001 — degrade to empty, never break ingestion
        logger.warning("PDF extraction failed: %s", exc)
        return ""


def extract_docx(content: bytes) -> str:
    try:
        from docx import Document

        doc = Document(io.BytesIO(content))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception as exc:  # noqa: BLE001
        logger.warning("DOCX extraction failed: %s", exc)
        return ""


def extract_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation

        slides: list[str] = []
        prs = Presentation(io.BytesIO(content))
        for slide in prs.slides:
            shapes: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shapes.append(shape.text)
            if shapes:
                slides.append("\n".join(shapes))
        return "\n\n---\n\n".join(slides)
    except Exception as exc:  # noqa: BLE001
        logger.warning("PPTX extraction failed: %s", exc)
        return ""


# --- segment-aware variants (preserve page/slide labels for citation traceability) -----------


def segments_pdf(content: bytes) -> list[tuple[str, str]]:
    """Return ``[("p.N", page_text), ...]`` — one entry per non-empty PDF page."""
    try:
        import pdfplumber

        segments: list[tuple[str, str]] = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text()
                if text and text.strip():
                    segments.append((f"p.{i}", text))
        return segments
    except Exception as exc:  # noqa: BLE001
        logger.warning("PDF segment extraction failed: %s", exc)
        return []


def segments_docx(content: bytes) -> list[tuple[str, str]]:
    """Return ``[("¶", body_text)]`` — DOCX has no reliable page boundary, so one body segment.

    Word page breaks are rendered at layout time, not stored in the XML, so a faithful page
    label isn't recoverable without rendering. One body segment keeps the label honest.
    """
    try:
        from docx import Document

        doc = Document(io.BytesIO(content))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        body = "\n".join(parts)
        return [("body", body)] if body.strip() else []
    except Exception as exc:  # noqa: BLE001
        logger.warning("DOCX segment extraction failed: %s", exc)
        return []


def segments_pptx(content: bytes) -> list[tuple[str, str]]:
    """Return ``[("slide N", slide_text), ...]`` — one entry per non-empty slide."""
    try:
        from pptx import Presentation

        segments: list[tuple[str, str]] = []
        prs = Presentation(io.BytesIO(content))
        for i, slide in enumerate(prs.slides, start=1):
            shapes = [
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if shapes:
                segments.append((f"slide {i}", "\n".join(shapes)))
        return segments
    except Exception as exc:  # noqa: BLE001
        logger.warning("PPTX segment extraction failed: %s", exc)
        return []
